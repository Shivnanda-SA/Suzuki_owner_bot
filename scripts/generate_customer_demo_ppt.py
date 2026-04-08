"""Generate Suzuki Owner Support PoC customer demo deck (python-pptx)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    sub.text = subtitle


def _bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)


def _two_col_bullets(prs: Presentation, title: str, left: list[str], right: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.75))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(32)
    tx.text_frame.paragraphs[0].font.bold = True

    def add_box(x: float, lines: list[str]) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(1.35), Inches(4.4), Inches(5.5))
        tf = box.text_frame
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(16)

    add_box(0.5, left)
    add_box(5.1, right)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(
        prs,
        "Suzuki Owner Support",
        "AI-powered post-purchase assistance — Proof of Concept\nCustomer demo deck",
    )

    _bullet_slide(
        prs,
        "Agenda",
        [
            "Context: why this PoC exists",
            "What we built (solution in one slide)",
            "How it works (architecture & trust)",
            "Live demo walkthrough",
            "Deployment & path to production",
            "Q&A",
        ],
    )

    _bullet_slide(
        prs,
        "The challenge",
        [
            "Owners need quick, accurate answers after purchase: service intervals, warranty, manuals, service centers, campaigns.",
            "Information is spread across manuals, websites, and dealer channels — hard to search consistently.",
            "Support teams scale cost and response time; generic chatbots risk wrong answers on regulated topics.",
        ],
    )

    _bullet_slide(
        prs,
        "What this PoC is",
        [
            "A working prototype: conversational “Owner Support” for Suzuki Motorcycle India (scooters & motorcycles).",
            "Grounded answers: retrieval from an indexed knowledge base, not free-form guessing.",
            "Synthetic / representative seed content for demo — replace with crawled public content + legal review before production.",
        ],
    )

    _bullet_slide(
        prs,
        "Solution we are demonstrating",
        [
            "Web chat UI: vehicle model selection, natural-language questions, quick prompts (first service, warranty, centers, campaigns).",
            "Backend API: RAG pipeline — retrieve relevant chunks → generate concise answers with citations.",
            "Structured responses: detected intent, confidence, source excerpts, optional CTA (e.g. official site), handoff flag when human help is needed.",
            "Session + feedback hooks: PostgreSQL for chat sessions and user feedback for iteration.",
        ],
    )

    _two_col_bullets(
        prs,
        "Architecture (high level)",
        [
            "Next.js web app",
            "FastAPI REST API",
            "LangGraph: retrieve → generate",
            "ChromaDB: vector store (persisted volume)",
            "OpenAI: chat + embeddings",
            "PostgreSQL: sessions & feedback",
        ],
        [
            "Flow:",
            "1. User message + optional model hint",
            "2. Similarity search (top-k chunks)",
            "3. LLM answers only from context",
            "4. JSON: answer + citations + intent",
            "",
            "Deploy: Docker Compose locally; Render template for cloud.",
        ],
    )

    _bullet_slide(
        prs,
        "Trust & safety (built into the design)",
        [
            "System prompt restricts answers to retrieved context; avoids inventing warranty or service promises.",
            "Citations: document title, type, source URL, excerpt — transparent provenance.",
            "needs_handoff when escalation to authorized service / official support is appropriate.",
        ],
    )

    _bullet_slide(
        prs,
        "Suggested live demo (5–8 minutes)",
        [
            "Open the web app; select a model (e.g. Access 125 or Gixxer SF 250).",
            "Try a quick prompt: first service timing, warranty coverage, service center in a city, service campaigns.",
            "Show citations under the answer and explain they map to indexed content.",
            "Optional: show API docs at /docs and health check for ops stakeholders.",
        ],
    )

    _bullet_slide(
        prs,
        "From PoC to production (talk track)",
        [
            "Replace seed JSON with approved content pipeline (official pages, manuals PDFs as allowed).",
            "Legal / compliance review of answers and disclaimers; brand tone guidelines.",
            "Observability: logging, rate limits, admin reindex (secured), monitoring on Render or enterprise hosting.",
            "Integrate with CRM or dealer systems only where scope is agreed.",
        ],
    )

    _title_slide(
        prs,
        "Thank you",
        "Questions?\n\nReference: RUNBOOK.txt for local run, Docker, and Render deployment.",
    )

    return prs


def main() -> None:
    out = Path(__file__).resolve().parent.parent / "Suzuki_Owner_Support_PoC_Customer_Demo.pptx"
    prs = build()
    prs.save(out)
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
